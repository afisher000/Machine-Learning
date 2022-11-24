import collections.abc #needed for pptx import
from pptx import Presentation
import os

class Notes():
    def __init__(self, main_gui):
        self.main_gui = main_gui
        self.notes_file = 'slide_notes.pptx'
        
    def save_notes(self, title, text, canvas):
        path = os.path.join(self.main_gui.directory, self.notes_file)
        ppt = Presentation(path) if os.path.exists(path) else Presentation()
        
        png_file = 'temp.png'
        canvas.print_figure(png_file, bbox_inches='tight')

        # Add slide with text and picture
        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
        slide.shapes[0].text_frame.paragraphs[0].text = ' ' if title=='' else title
        slide.shapes[2].text_frame.paragraphs[0].text = text
        
        # Add picture with no cropping
        pic = slide.placeholders[1].insert_picture(png_file)
        pic.crop_top, pic.crop_left, pic.crop_bottom, pic.crop_right = 0,0,0,0
        
        # Save powerpoint and delete temporary figure file
        ppt.save(path)
        os.remove(png_file)
        return