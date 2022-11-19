# %%
import sys
sys.path.append('../Datasets')
from ClassifierAnalysis import ClassifierModels
from RegressorAnalysis import RegressorModels
from Pipelines import Imputation, Scaling
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import PolynomialFeatures

import collections.abc #needed for pptx import
from pptx import Presentation
import sys
import os
import pandas as pd
from pandas.api.types import is_string_dtype, is_numeric_dtype
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt

# PyQt5 imports
from PyQt5 import QtCore as qtc
from PyQt5 import QtWidgets as qtw
from PyQt5 import QtGui as qtg
from PyQt5.Qt import Qt as qt
from PyQt5 import uic

# Figure Imports
from matplotlib.backends.backend_qt5agg import FigureCanvas, NavigationToolbar2QT as NavigationToolbar
from matplotlib.figure import Figure
import matplotlib.gridspec as gridspec
from matplotlib import rcParams
rcParams.update({'figure.autolayout': True}) # to automatically fit long axis ticklabels

from sklearn.base import BaseEstimator, TransformerMixin
class FeatureEngineeringTransformations(BaseEstimator, TransformerMixin):
    def __init__(self, gui, file):
        self.gui = gui
        self.file = file

    def fit(self, df, y=None):
        return self

    def transform(self, df, y=None):
        with open(self.file,'r') as f:
            code = f.read()
        try:
            exec(code)
        except Exception as e:
            qtw.QMessageBox.critical(self.gui, 'Feature Engineering Error', f'Raised following error: {e}')
        return df

class KeepSelectedFeatures(BaseEstimator, TransformerMixin):
    def __init__(self, features):
        self.features = features

    def fit(self, df, y=None):
        return self

    def transform(self, df, y=None):
        return df[self.features]
    
class FeatureEngr():
    def __init__(self, file_manager):
        self.file_manager = file_manager
        self.data_file = 'feature_engineering.csv'
        self.featengr_file = 'feature_engineering_code.txt'
        self.settings_file = 'feature_engineering_settings.csv'
        
    def new_feature(self, transform, feature=None, map_str=None, code=None):
        if transform=='log1p':
            self.data['log'+feature] = np.log1p(self.data[feature])
            self.settings.loc['log'+feature] = ['log'+feature, False, 'num']
            self.save_featengr_code( 
                feature,
                f'df["log{feature}"]=np.log1p(df["{feature}"])'
            )

        elif transform == 'bin':
            self.data['bin'+feature] = pd.qcut(self.data[feature], 10)
            self.settings.loc['bin'+feature] = ['bin'+feature, False, 'cat']
            self.save_featengr_code( 
                feature,
                f'df["bin{feature}"]=pd.qcut(df["{feature}"], 10)'
            )

        elif transform == 'dummies':
            self.data = pd.concat([ 
                self.data, 
                pd.get_dummies(self.data[feature], prefix=feature)], axis=1
            )
            if 'Null' in self.data[feature].cat.categories:
                self.data = self.data.drop(columns=[feature+'_Null'])
                
            self.save_featengr_code( 
                feature,
                (
                    f'df=pd.concat([df, pd.get_dummies(df["{feature}"], prefix="{feature}")], axis=1)' + '\n'
                    f'if "Null" in df["{feature}"].cat.categories: df=df.drop(columns=["{feature}_Null"])'
                )
            )

            for cat in self.data[feature].unique():
                if cat!='Null':
                    self.settings.loc[f'{feature}_{cat}'] = [f'{feature}_{cat}', False, 'pnum']

        elif transform == 'drop':
            self.settings.loc[feature, 'isdropped'] = True

        elif transform == 'undrop':
            self.settings.loc[feature, 'isdropped'] = False

        elif transform == 'encode':
            try:
                encode_map = {}
                for entry in map_str.split(','):
                    key, value = entry.split('=')
                    if key.strip() not in self.data[feature].values:
                        raise ValueError(f'Key "{key}" is invalid for given feature')
                    encode_map[key.strip()] = float(value.strip())
            except Exception as e:
                self.file_manager.data_gui.raise_error('Invalid Map String', f'{e}')
                return 1

            try:
                self.data['enc'+feature] = self.data[feature].map(encode_map).astype('float')
            except Exception as e:
                self.file_manager.data_gui.raise_error('Encoding Error', f'{e}')
                return 1

            self.save_featengr_code( 
                feature, 
                f'df["enc_{feature}"] = df["{feature}"].map({encode_map}).astype("float")'
            )
            self.settings.loc['enc_'+feature] = ['enc_'+feature, False, 'pnum']

        elif transform == 'code':
            df = self.data
            old_features = df.columns
            try:
                exec(code)
            except Exception as e:
                self.file_manager.data_gui.raise_error('Code Error', f'{e}')
                return 1
            self.save_featengr_code(None, code)

            # Make parsings feature type a function?
            for feature in (set(df.columns)-set(old_features)):
                self.settings.loc[feature] = [feature, False, 'none']
                if df[feature].dtype == 'object':
                    df[feature] = df[feature].fillna('Null')
                    df[feature] = pd.Categorical(df[feature], np.sort(df[feature].unique()))
                    self.settings.loc[feature, 'feature_type'] = 'cat'
                else:
                    if df[feature].nunique()<11: #TODO make 11 and attribute? Used multiple locations
                        self.settings.loc[feature, 'feature_type'] = 'pnum'
                    else:
                        self.settings.loc[feature, 'feature_type'] = 'num'
    
        self.data.to_csv(os.path.join(self.file_manager.directory, self.data_file), index=False)
        self.settings.to_csv(os.path.join(self.file_manager.directory, self.settings_file))
        return

    def save_featengr_code(self, feature, code):
        # TODO Check if feature is target feature...
        with open(os.path.join(self.file_manager.directory, self.featengr_file), 'a') as f:
            f.write(code)
            f.write('\n')

    def get_feature_types(self):
        return self.settings.feature_type
    
    def get_discrete_features(self):
        return self.data.columns[self.data.nunique()<0.05*len(self.data)].to_list()
        
    def get_features_by_type(self, feature_type):
        # Sort?
        feature_list = self.settings.name[self.settings.feature_type==feature_type].to_list()
        return sorted(feature_list)

    def get_features_by_isdropped(self, isdropped, prefix=False):
        # Sort?
        if prefix:
            full_name = self.settings.feature_type + '_' + self.settings.name
        else:
            full_name = self.settings.name
        feature_list = full_name[self.settings.isdropped==isdropped].to_list()
        return sorted(feature_list)

    def create_files(self):
        # Read and differentiate train data
        data = pd.read_csv(os.path.join(self.file_manager.directory, 'train.csv'))
        obj_features = data.select_dtypes('object').columns
        num_features = data.select_dtypes('number').columns
        is_pnum = data[num_features].nunique().lt(11)

        # Convert to object dtypes to Categoricals
        data[obj_features] = data[obj_features].fillna('Null')
        for feature in obj_features.to_series():
            data[feature] = pd.Categorical(data[feature], np.sort(data[feature].unique()))
        
        # Build featsettings
        settings = pd.DataFrame(index=data.columns)
        settings['name'] = data.columns
        settings['isdropped'] = False
        settings.loc[obj_features, 'feature_type'] = 'cat'
        settings.loc[num_features[is_pnum], 'feature_type'] = 'pnum'
        settings.loc[num_features[~is_pnum], 'feature_type'] = 'num'

        # Save to file
        data.to_csv(os.path.join(self.file_manager.directory, self.data_file), index=False)
        settings.to_csv(os.path.join(self.file_manager.directory, self.settings_file))
        return

    def load_files(self):
        # Read or create featengr
        data_path = os.path.join(self.file_manager.directory, self.data_file)
        settings_path = os.path.join(self.file_manager.directory, self.settings_file)

        if not os.path.exists(data_path):
            self.create_files()

        self.data = pd.read_csv(data_path)
        self.settings = pd.read_csv(settings_path, index_col=0)
        
        # Convert to object dtypes to Categoricals
        obj_features = self.data.select_dtypes('object').columns
        self.data[obj_features] = self.data[obj_features].fillna('Null')
        for feature in obj_features.to_series():
            self.data[feature] = pd.Categorical(self.data[feature], np.sort(self.data[feature].unique()))
        return

        
class Notes():
    def __init__(self, file_manager):
        self.file_manager = file_manager
        self.notes_file = 'slide_notes.pptx'
        
    def save_notes(self, title, text, canvas):
        path = os.path.join(self.file_manager.directory, self.notes_file)
        ppt = Presentation(path) if os.path.exists(path) else Presentation()
        
        png_file = 'temp.png'
        canvas.print_figure(png_file, bbox_inches='tight')

        # Add slide with text and picture
        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
        slide.shapes[0].text_frame.paragraphs[0].text = ' ' if title=='' else title
        slide.shapes[2].text_frame.paragraphs[0].text = text
        
        # Add picture with no cropping
        pic = slide.placeholders[1].insert_picture(png_file)
        pic.crop_top, pic.crop_left, pic.crop_bottom, pic.crop_right = 0,0,0,0
        
        # Save powerpoint and delete temporary figure file
        ppt.save(path)
        os.remove(png_file)
        return
        
class Pipe():
    def __init__(self, file_manager):
        self.file_manager = file_manager
        self.settings_file = 'pipe_settings.csv'
        self.widget_ptrs = None
        
    def load_files(self):
        # Create file if does not exists
        path = os.path.join(self.file_manager.directory, self.settings_file)
        if os.path.exists(path):
            self.settings = pd.read_csv(path, index_col=0)
        else:
            self.settings = pd.DataFrame(columns=['is_selected','scale_strat','impute_strat', 'impute_by'])
            
        
        # Update with new feature engineered features
        features = self.file_manager.featengr.data.columns
        for feature in features:
            if feature not in self.settings.index:
                self.settings.loc[feature] = [False, 'none','none','none']
        self.settings.to_csv(path, index=True)
        
        
        # Create widget pointer dataframe
        self.widget_ptrs = pd.DataFrame(index=features, columns=['checkbox','scaling_menu','imputing_menu','imputeby_menu','fillna_input'])
        return
    
    def update_widgets(self):
        for feature, row in self.widget_ptrs.iterrows():
            is_selected, scale_strat, impute_strat, impute_by = self.settings.loc[feature]
            
            row.checkbox.setChecked(is_selected)
            row.scaling_menu.setCurrentText(scale_strat)
            row.imputing_menu.setCurrentText(impute_strat)
            if impute_strat=='bycategory':
                row.imputeby_menu.setCurrentText(impute_by)
            elif impute_strat=='byvalue':
                row.fillna_input.setText(impute_by)
        return
    
    def get_selected_features(self):
        return self.settings.index[self.settings.is_selected]
    
    def save_settings(self):
        
        for feature, row in self.widget_ptrs.iterrows():
            is_selected = row.checkbox.isChecked()
            scale_strat = row.scaling_menu.currentText()
            impute_strat = row.imputing_menu.currentText()
            if impute_strat=='bycategory':
                impute_by = row.imputeby_menu.currentText()
            elif impute_strat=='byvalue':
                impute_by = row.fillna_input.text()
            else:
                impute_by = 'none'
            self.settings.loc[feature] = [is_selected, scale_strat, impute_strat, impute_by]
        
        path = os.path.join(self.file_manager.directory, self.settings_file)
        self.settings.to_csv(path)

    def get_scaling_dict(self):
        data = {}
        strategies = self.settings[self.settings.is_selected].scale_strat
        for feature, strat in strategies.items():
            if strat=='none':
                continue
            if strat in data:
                data[strat].append(feature)
            else:
                data[strat] = [feature]
        return data  
        
    def get_imputation_dict(self):
        # Build dictionary for imputation implementation
        data = {}
        strategies = self.settings[self.settings.is_selected].impute_strat
        for feature, strat in strategies.items():
            if strat == 'none':
                continue
            
            # Entry in dictinoary depends on strategy
            if strat in ['bycategory', 'byvalue']:
                by_value = self.settings.impute_by[feature]
                entry = (feature, by_value)
            else:
                entry = feature
                
            if strat in data:
                data[strat].append(entry)
            else:
                data[strat] = [entry]
        return data
    
    
class FileManager():
    def __init__(self):
        self.data_gui = None
        self.model_gui = None
        self.featengr = FeatureEngr(self)
        self.notes = Notes(self)
        self.pipe = Pipe(self)
        self.directory = '.' # initialize to current directory
        return


    def load_directory(self, directory):
        self.directory = directory
        
        # Get target and ID features
        train_df = pd.read_csv(os.path.join(directory, 'train.csv'))
        test_df = pd.read_csv(os.path.join(directory, 'test.csv'))
        self.id_feature = test_df.columns[0] #Can be done better
        self.target_feature = (set(train_df.columns) - set(test_df.columns)).pop()
        
        # Load data files
        self.featengr.load_files()
        self.pipe.load_files()


mw_Ui, mw_Base = uic.loadUiType('master_window.ui')
class MainGUI(mw_Base, mw_Ui):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.show()

        # Import Layout
        self.setupUi(self)
        self.file_manager = FileManager()
        

        # Connect signals and slots
        self.browseDirectory_button.clicked.connect(self.select_directory)
        self.featureEngineering_button.clicked.connect(self.show_data_gui)
        self.modelAnalysis_button.clicked.connect(self.show_model_gui)

    def select_directory(self):
        directory = str(qtw.QFileDialog.getExistingDirectory(self, 'Select Directory'))
        if ('train.csv' not in os.listdir(directory)) or ('test.csv' not in os.listdir(directory)):
            qtw.QMessageBox.critical(self, 'Invalid Directory', 'Directory must include "test.csv" and "train.csv" files.')
            return
        
        self.directory_label.setText(os.path.relpath(directory))
        self.file_manager.load_directory(directory)
        
    def show_data_gui(self):
        self.data_gui = DataGUI(self)
        
    def show_model_gui(self):
        self.model_gui = ModelGUI(self)
        


mw_Ui, mw_Base = uic.loadUiType('model_analysis_gui.ui')
class ModelGUI(mw_Base, mw_Ui):
    def __init__(self, main_gui, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.show()
        
        # Import Layout
        self.setupUi(self)
        self.file_manager = main_gui.file_manager
        self.file_manager.pipe.load_files()
        self.file_manager.model_gui = self
        self.main_gui = main_gui
        self.figure = GUI_Figure(self, self.plot_layout)
        self.populate_pipeline_widgets()
        
        # Signals and Slots
        self.applypipeline_button.clicked.connect(self.apply_pipeline)
        
        
    def apply_pipeline(self):
        # TODO - something to show pipeline was updated
        self.file_manager.pipe.save_settings()
        
        train = pd.read_csv(os.path.join(self.file_manager.directory, 'train.csv'))
        test = pd.read_csv(os.path.join(self.file_manager.directory, 'test.csv'))
        id_feature = self.file_manager.id_feature
        target_feature = self.file_manager.target_feature
        y = train[target_feature]
        
        scaling_dict = self.file_manager.pipe.get_scaling_dict()
        imputation_dict = self.file_manager.pipe.get_imputation_dict()
        featengr_path = os.path.join(self.file_manager.directory, self.file_manager.featengr.featengr_file)
        selected_features = self.file_manager.pipe.get_selected_features()
        transformations = [ 
            ('Feature Engineering', FeatureEngineeringTransformations(self, featengr_path)),
            ('Imputation', Imputation(imputation_dict)),
            ('Scaling', Scaling(scaling_dict)),
            ('DropFeatures', KeepSelectedFeatures(selected_features))
        ]
        
        # Polynomial Features
        if self.polyfeatures_checkbox.isChecked():
            polydeg = self.polydeg_spinbox.value()
            interaction_only = self.polyinteraction_checkbox.isChecked()
            transformations.append(('PolyFeatures', PolynomialFeatures(polydeg, interaction_only=interaction_only)))

        pipeline = Pipeline(transformations)
        X = pipeline.fit_transform(train)
        X_test = pipeline.transform(test)
        
        
    def populate_pipeline_widgets(self):
        grid = self.pipeline.layout()
        
        # Get discrete features for bycategory menu
        discrete_features_plus_none = self.file_manager.featengr.get_discrete_features()
        discrete_features_plus_none.insert(0, 'none')

        
        # Create title row
        titles = ['Feature','Scaling','Imputing','Category','Value']
        for icol, title in enumerate(titles):
            grid.addWidget(qtw.QLabel(title, self), 0, icol)

        # Populate by row
        for irow, (feature, row) in enumerate(self.file_manager.pipe.settings.iterrows()):

            # Create widgets
            box = qtw.QCheckBox(feature, self)
            scaling_menu = qtw.QComboBox(self)
            scaling_menu.addItems(['none','minmax','standard','quartile'])
            imputing_menu = qtw.QComboBox(self)
            imputing_menu.addItems(['none','mean','median','most_frequent','bycategory','byvalue'])
            imputeby_menu = qtw.QComboBox(self)
            imputeby_menu.addItems(discrete_features_plus_none) 
            fillna_input = qtw.QLineEdit(self)

            # Link pointers in table
            self.file_manager.pipe.widget_ptrs.loc[feature] = [box, scaling_menu, imputing_menu, imputeby_menu, fillna_input]

            # Populate row
            grid.addWidget(box, irow+1, 0)
            grid.addWidget(scaling_menu, irow+1, 1)
            grid.addWidget(imputing_menu, irow+1, 2)
            grid.addWidget(imputeby_menu, irow+1, 3)
            grid.addWidget(fillna_input, irow+1, 4)

            
        grid.setRowStretch(grid.rowCount(), 1)
        self.file_manager.pipe.update_widgets()
        return   

mw_Ui, mw_Base = uic.loadUiType('data_analysis_gui.ui')
class DataGUI(mw_Base, mw_Ui):
    def __init__(self, main_gui, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.show()
        
        # Import Layout
        self.setupUi(self)
        self.file_manager = main_gui.file_manager
        self.file_manager.data_gui = self
        self.main_gui = main_gui
        self.figure = GUI_Figure(self, self.plot_layout)
        self.update_menus()
        self.initialize_fixed_menus()

        # Connect signals and slots
        self.log_button.clicked.connect(self.log_transform_feature)
        self.bin_button.clicked.connect(self.bin_transform_feature)
        self.create_button.clicked.connect(self.create_feature_from_code)
        self.dummies_button.clicked.connect(self.dummies_transform_feature)
        self.encode_button.clicked.connect(self.encode_transform_feature)
        self.drop_button.clicked.connect(self.drop_feature)
        self.undrop_button.clicked.connect(self.undrop_feature)
        self.notes_button.clicked.connect(self.save_notes)
        self.update_button.clicked.connect(self.update_figure)
        self.customplot_button.clicked.connect(self.custom_plot)

    def initialize_fixed_menus(self):
        # Setup aggfcn_combobox
        agg_fcn_options = {'none':None, 'max':np.max, 'min':np.min, 'mean':np.mean, 'median':np.median, 'sum':np.sum}
        for label, fcn in agg_fcn_options.items():
            self.aggfcn_combobox.addItem(label, fcn)
            
        # Setup customplot_combobox
        self.customplot_combobox.addItems(['none','correlation', 'null percentage', 'largest category'])
        
        
        
    def log_transform_feature(self):
        feature = self.log_combobox.currentText()
        if feature == 'none':
            return
        self.file_manager.featengr.new_feature('log1p', feature=feature)
        self.log_combobox.setCurrentText('none')
        self.update_menus()

    def bin_transform_feature(self):
        feature = self.bin_combobox.currentText()
        if feature == 'none':
            return
        self.file_manager.featengr.new_feature('bin', feature=feature)
        self.bin_combobox.setCurrentText('none')
        self.update_menus()

    def drop_feature(self):
        prefix, feature = self.drop_combobox.currentText().split('_')
        if feature == 'none':
            return
        self.file_manager.featengr.new_feature('drop', feature=feature)
        self.drop_combobox.setCurrentText('none')
        self.update_menus()

    def dummies_transform_feature(self):
        feature = self.dummies_combobox.currentText()
        if feature == 'none':
            return 
        self.file_manager.featengr.new_feature('dummies', feature=feature)
        self.dummies_combobox.setCurrentText('none')
        self.update_menus()

    def undrop_feature(self):
        feature = self.undrop_combobox.currentText()
        if feature == '':
            return
        self.file_manager.featengr.new_feature('undrop', feature=feature)
        self.update_menus()

    def create_feature_from_code(self):
        code = self.create_lineedit.text()
        if not self.file_manager.featengr.new_feature('code', code=code):
            self.create_lineedit.setText('')
            self.update_menus()

    def encode_transform_feature(self):
        map_str = self.encode_lineedit.text()
        feature = self.encode_combobox.currentText()
        if feature == 'none':
            return

        if not self.file_manager.featengr.new_feature('encode', feature=feature, map_str=map_str):
            self.encode_combobox.setCurrentText('none')
            self.encode_lineedit.setText('')
            self.update_menus()

    def raise_error(self, title, text):
        qtw.QMessageBox.critical(self, title, text)
        return

    def update_menus(self):
        dropped_features = self.file_manager.featengr.get_features_by_isdropped(True)
        undropped_features = self.file_manager.featengr.get_features_by_isdropped(False, prefix=True)

        undropped_features_plus_none = undropped_features
        undropped_features_plus_none.insert(0, 'none')

        num_features_plus_none = self.file_manager.featengr.get_features_by_type('num')
        num_features_plus_none.insert(0, 'none')

        cat_features_plus_none = self.file_manager.featengr.get_features_by_type('cat')
        cat_features_plus_none.insert(0, 'none')

        # Specify options for each combobox
        options_dict = { 
            'x':undropped_features_plus_none,
            'y':undropped_features_plus_none,
            'hue':undropped_features_plus_none,
            'size':undropped_features_plus_none,
            'style':undropped_features_plus_none,
            'log':num_features_plus_none,
            'bin':num_features_plus_none,
            'dummies':cat_features_plus_none,
            'encode':cat_features_plus_none,
            'drop':undropped_features_plus_none,
            'undrop':dropped_features
        }

        # Update combobox items
        for menu, options in options_dict.items():
            combobox = getattr(self, menu+'_combobox')
            current_text = combobox.currentText()
            combobox.clear()
            combobox.addItems(options)
            combobox.setCurrentText(current_text)

    def save_notes(self):
        ''' Write comments and figure to notes powerpoint'''
        title = self.title_lineedit.text()
        text = self.text_textedit.toPlainText()
        canvas = self.figure.canvas
        self.file_manager.notes.save_notes(title, text, canvas)
        
        # Reset GUI inputs
        self.title_lineedit.setText('')
        self.text_textedit.setText('')
        
    def parse_combobox_text(self, comboboxes):
        is_not_list = not isinstance(comboboxes, list)
        if is_not_list:
            comboboxes = [comboboxes]
        output = []
        for combobox in comboboxes:
            text = combobox.currentText()
            if text=='none':
                output.append(None)
            else:
                prefix, feature = text.split('_')
                output.append(feature)
                
        if is_not_list:
            output = output[0]
        return output
        
    def update_figure(self):
        x, y, hue, size, style = self.parse_combobox_text( 
            [self.x_combobox, self.y_combobox, self.hue_combobox, self.size_combobox, self.style_combobox]
        )
        agg_fcn = self.aggfcn_combobox.currentData()
        
        data = self.file_manager.featengr.data
        num_features = self.file_manager.featengr.get_features_by_type('num')
        pnum_features = self.file_manager.featengr.get_features_by_type('pnum')
        cat_features = self.file_manager.featengr.get_features_by_type('cat')
        
        if x is None:
            qtw.QMessageBox.warning(self, 'No X specified', 'The X feature must be specified to plot')
            self.figure.reset_figure(ncols=1)
            self.figure.draw()
            return 
        
        if y is None:
            if x in num_features:
                self.figure.histplot(data, x, hue)
                return
            else:
                self.figure.histplot_and_piechart(data, x, hue)
                return
        
        if x in num_features:
            if y in num_features:
                self.figure.scatterplot(data, x, y, hue, size, style)
            elif y in pnum_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, switch_axes=True
                )
            elif y in cat_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, switch_axes=True, median_ordering=True
                )
        elif x in pnum_features:
            if y in num_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, agg_fcn=agg_fcn 
                )
            elif y in pnum_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, agg_fcn=agg_fcn 
                )
            elif y in cat_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, switch_axes=True, median_ordering=True 
                )
        elif x in cat_features:
            if y in num_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, median_ordering=True 
                )
            elif y in pnum_features:
                self.figure.boxenplot_and_piechart( 
                    data, x, y, hue, agg_fcn = agg_fcn
                )     
            elif y in cat_features:
                self.figure.histplot_and_piechart( 
                    data, x, hue=y, stat='count'
                )
        
    def custom_plot(self):
        plot_type = self.customplot_combobox.currentText()
        if plot_type=='none':
            return
        
        if self.file_manager.target_feature in self.file_manager.featengr.get_features_by_type('cat') and plot_type=='correlation':
            qtw.QMessageBox.warning(
                self, 'Invalid Target Type', 'Correlations are not valid for categorical targets.'
            )
            return
        n_features = self.customplot_spinbox.value()
        data = self.file_manager.featengr.data
        undropped_features = self.file_manager.featengr.get_features_by_isdropped(False)
        self.figure.custom_plot(data[undropped_features], plot_type, self.file_manager.target_feature, n_features)
        return
        
class GUI_Figure():
    def __init__(self, GUI, layout):
        self.GUI = GUI
        self.canvas = FigureCanvas(Figure(figsize=(3,3)))
        layout.addWidget(NavigationToolbar(self.canvas, self.GUI))
        layout.addWidget(self.canvas)
    
    def reset_figure(self, ncols):
        self.canvas.figure.clear()
        self.ax = self.canvas.figure.subplots(ncols=ncols)
        return
    
    def custom_plot(self, data, plot_type, target_feature, n_features):
        self.reset_figure(ncols=1)
        if plot_type == 'correlation':
            ( 
                data.corrwith(data[target_feature]) #, numeric_only=True
                .abs().sort_values(ascending=False)[1:n_features+1]
                .plot.bar(xlabel='Feature', ylabel='Correlation', title=f'Target={target_feature}', ax=self.ax)
            )
        elif plot_type == 'null percentage':
            ( 
                (data.isna().sum()+data.apply(lambda x: sum(x=='Null')))
                .sort_values(ascending=False).divide(len(data))[:n_features]
                .plot.bar(xlabel='Feature', ylabel='Null Percentage', ax=self.ax)
            )
        elif plot_type == 'largest category':
            ( 
                data.apply(lambda x: x.value_counts().iloc[0]/len(x)).
                sort_values(ascending=False)[:n_features]
                .plot.bar(xlabel='Feature', ylabel='Largest Cat. Fraction', ax=self.ax)
            )
        self.draw()
        
    def boxenplot_and_piechart(self, data, x, y, hue, agg_fcn=None, switch_axes=False, median_ordering=False):
        self.reset_figure(ncols=2)
        if switch_axes:
            x, y = y, x
            
        if median_ordering:
            data = data.copy()
            median_ordering = data.groupby(by=x)[y].median().sort_values().index.tolist()
            data[x] = data[x].cat.reorder_categories(median_ordering)
            
        if agg_fcn is None:
            sns.boxenplot(data=data, x=x, y=y, hue=hue, ax=self.ax[0])
        else:
            sns.barplot(data=data.dropna(subset=[x,y]), x=x, y=y, hue=hue, estimator=agg_fcn, ax=self.ax[0])

        self.ax[0].tick_params(axis='x', rotation=90)
        data[x].value_counts(dropna=False).sort_index().plot.pie(autopct='%.0f%%', ax=self.ax[1])
        self.draw()
        return 

    
    def histplot_and_piechart(self, data, x, hue, stat=None):
        # Add check for number of unique in x?
        self.reset_figure(ncols=2)
        sns.histplot(data=data, x=x, hue=hue, stat='count', multiple='stack', kde=False, ax=self.ax[0])
        self.ax[0].tick_params(axis='x', rotation=90)
        data[x].value_counts(dropna=False).sort_index().plot.pie(autopct='%.0f%%', ax=self.ax[1])
        self.draw()
        return
        
    def histplot(self, data, x, hue):
        self.reset_figure(ncols=1)
        sns.histplot(data=data, x=x, hue=hue, multiple='stack', kde=True, ax=self.ax)
        self.draw()
        return
    
    def scatterplot(self, data, x, y, hue, size, style):
        self.reset_figure(ncols=1)
        sns.scatterplot(data=data, x=x, y=y, hue=hue, size=size, style=style, 
                        ax=self.ax)
        self.draw()
        return
    
        
        
    def draw(self):
        self.canvas.draw()

    

if __name__ =='__main__':
    app = qtw.QApplication(sys.argv)
    gui = MainGUI()
    app.exec()
# %%
