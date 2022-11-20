# %%  -*- coding: utf-8 -*-

"""
Created on Mon Oct 24 08:40:32 2022

@author: afisher
"""

import collections.abc #needed for pptx import
from pptx import Presentation
import sys
import os
import pandas as pd
from pandas.api.types import is_string_dtype, is_numeric_dtype
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt

# PyQt5 imports
from PyQt5 import QtCore as qtc
from PyQt5 import QtWidgets as qtw
from PyQt5 import QtGui as qtg
from PyQt5.Qt import Qt as qt
from PyQt5 import uic

# Figure Imports
from matplotlib.backends.backend_qt5agg import FigureCanvas, NavigationToolbar2QT as NavigationToolbar
from matplotlib.figure import Figure
import matplotlib.gridspec as gridspec
from matplotlib import rcParams
rcParams.update({'figure.autolayout': True}) # to automatically fit long axis ticklabels

## To implement
# Filtering options (only data from a given category or numeric range)

mw_Ui, mw_Base = uic.loadUiType('data_analysis_gui.ui')
class Main_Window(mw_Base, mw_Ui):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.show()

        # Attributes
        self.file_directory = None

        # Import Layout
        self.setupUi(self)

        # Add to GUI
            # Initialize figure
        self.canvas = FigureCanvas(Figure(figsize=(3,3)))
        self.plot_layout.addWidget(NavigationToolbar(self.canvas, self))
        self.plot_layout.addWidget(self.canvas)

            # Add options to aggfcn combobox
        aggregate_fcn_options = {' none':None, ' max':np.max, ' min':np.min, ' mean':np.mean, ' median':np.median, ' sum':np.sum}
        for label, fcn in aggregate_fcn_options.items():
            self.aggregate_combobox.addItem(label, fcn)

        # Connect signals and slots
        self.browse_button.clicked.connect(self.browse_file)
        self.update_button.clicked.connect(self.update_figure)
        self.drop_button.clicked.connect(self.drop_feature)
        self.undrop_button.clicked.connect(self.undrop_feature)
        self.corr_button.clicked.connect(self.plot_correlations)
        self.nullpct_button.clicked.connect(self.plot_nullpct)
        self.log_button.clicked.connect(self.log_transform_feature)
        self.bin_button.clicked.connect(self.bin_transform_feature)
        self.saveas_button.clicked.connect(self.save_as)
        self.save_button.clicked.connect(self.save)
        self.notes_button.clicked.connect(self.save_notes_to_pptx)
        self.create_button.clicked.connect(self.create_feature_from_code)
        self.dummies_button.clicked.connect(self.dummies_transform_feature)
        self.encode_button.clicked.connect(self.encode_transform_feature)



    def browse_file(self):
        # Browse for file
        abs_filename, _ = qtw.QFileDialog.getOpenFileName()
        if self.abs_filename=='':
            return

        # Save file_directory, dataframe, and update Save LineEdit
        self.file_directory = os.path.dirname(abs_filename)
        self.save_lineedit.setText(os.path.relpath(abs_filename))
        self.df = pd.read_csv(abs_filename)

        # Identify dropped features by suffix and remove suffix
        self.dropped_features = [c.removesuffix('_dropped') 
            for c in self.df.columns if c.endswith('_dropped')]
        self.df = self.df.rename(columns={c:c.removesuffix('_dropped') 
            for c in self.df.columns if c.endswith('_dropped')})

        # Categorize undropped features as numerical (num), pseudonumerical (pnum), and cateogrical (cat)
        valid_features = self.df.drop(columns=self.dropped_features).columns.tolist()

        # Save 'object' dtypes as categoricals
        obj_dtype_features = self.df[valid_features].select_dtypes('object').columns
        self.df[obj_dtype_features] = self.df[obj_dtype_features].fillna('Null')
        self.cat_features = obj_dtype_features.to_series()
        for cat_feature in self.cat_features:
            self.df[cat_feature] = pd.Categorical(self.df[cat_feature], np.sort(self.df[cat_feature].unique()))

        # numerical and pseudonumeric
        num_dtype_features = self.df.select_dtypes('number').columns
        is_pseudo = self.df[num_dtype_features].nunique().lt(11)
        self.pnum_features = num_dtype_features[is_pseudo].to_series()
        self.num_features = num_dtype_features[~is_pseudo].to_series()

        # Update combobox menus
        self.update_menus()

    def save(self):
        # Handle path input
        rel_path = self.save_lineedit.text()
        if rel_path=='':
            qtw.QMessageBox.warning(self, 'No Filename', 'Add a filename to save to.')
            return
        if not rel_path.endswith('.csv'):
            rel_path += '.csv'

        
        df_to_save = self.df.rename(columns={c:c+'_dropped' for c in self.dropped_features})
        df_to_save.to_csv(rel_path, index=False)

    def save_as(self):
        if not hasattr(self, 'df'):
            qtw.QMessageBox.warning(self, 'Missing Data', 'No data available to save.')
            return

        # Handle filename
        filename = self.saveas_lineedit.text()
        if filename=='':
            qtw.QMessageBox.warning(self, 'No Filename', 'Add a filename to save to.')
            return
        if not filename.endswith('.csv'):
            filename += '.csv'
        full_path = os.path.join(self.file_directory, filename)

        # Encode dropped columns with '_dropped' suffix and save to file
        df_to_save = self.df.rename(columns={c:c+'_dropped' for c in self.dropped_features})
        df_to_save.to_csv(full_path, index=False)

        # Move file name from saveas lineedit to save lineedit
        self.saveas_lineedit.setText('')
        self.save_lineedit.setText(os.path.relpath(full_path))

    def reset_figure(self, ncols):
        '''Reset canvas figure'''
        self.canvas.figure.clear()
        self.ax = self.canvas.figure.subplots(ncols=ncols)
        return

    def update_figure(self):
        '''Update figure type depending on feature types. '''
        df = self.df
        
        def parse_combobox_text(combobox):
            prefix, feature = combobox.currentText().split(' ')
            return None if feature=='none' else feature

        def boxen_and_pie_plot(data, x, y, hue=None, agg_fcn=None):
            self.reset_figure(ncols=2)
            if agg_fcn is None:
                sns.boxenplot(data=data, x=x, y=y, hue=hue, ax=self.ax[0])
            else:
                sns.barplot(data=data.dropna(subset=[x,y]), x=x, y=y, hue=hue, estimator=agg_fcn, ax=self.ax[0], errorbar=None)
            self.ax[0].tick_params(axis='x', rotation=90)
            df[x].value_counts(dropna=False).sort_index().plot.pie(autopct='%.0f%%', ax=self.ax[1])

        x = parse_combobox_text(self.x_combobox)
        y = parse_combobox_text(self.y_combobox)
        hue = parse_combobox_text(self.hue_combobox)
        size = parse_combobox_text(self.size_combobox)
        style = parse_combobox_text(self.style_combobox)
        agg_fcn = self.aggregate_combobox.currentData()

        # Check that x is specified
        if x is None: 
            qtw.QMessageBox.warning(self, 'No X specified', 'The X feature must be specified to plot')
            self.reset_figure(ncols=1)
            self.canvas.draw()
            return

        # X input only
        if y is None:
            if x in self.num_features:
                self.reset_figure(ncols=1)
                sns.histplot(data=df, x=x, hue=hue, multiple='stack', kde=True, ax=self.ax)
            else:
                if df[x].nunique()>30:
                    qtw.QMessageBox.warning(self, 'Axis Overflow', 'Too many unique entried to plot.')
                    self.reset_figure(ncols=1)
                    self.canvas.draw()
                    return

                self.reset_figure(ncols=2)
                sns.histplot(data=df, x=x, hue=hue, multiple='stack', kde=False, ax=self.ax[0])
                self.ax[0].tick_params(axis='x', rotation=90)
                self.df[x].value_counts(dropna=False).sort_index().plot.pie(autopct='%.0f%%', ax=self.ax[1])

        # X-Y Inputs
        if (x in self.num_features):
            if (y in self.num_features):
                self.reset_figure(ncols=1)
                sns.scatterplot(data=df, x=x, y=y, hue=hue, size=size, style=style, ax=self.ax)
                # Add regression?
            elif (y in self.pnum_features):
                x, y = y, x #Switch pnum to x axis
                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue)

            elif (y in self.cat_features):
                x, y = y, x #Switch category to x axis
                # Sort according to median of numeric
                cat_ordering = df[x].cat.categories.tolist()
                median_ordering = df.groupby(by=x)[y].median().sort_values().index.tolist()
                df[x] = df[x].cat.reorder_categories(median_ordering)

                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue)

                # Revert to categorical ordering
                df[x] = df[x].cat.reorder_categories(cat_ordering)

        elif (x in self.pnum_features):
            if (y in self.num_features):
                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue, agg_fcn=agg_fcn)
            elif (y in self.pnum_features):
                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue, agg_fcn=agg_fcn)
            elif (y in self.cat_features):
                x, y = y, x #Switch category to x axis
                # Sort according to median of numeric
                cat_ordering = df[x].cat.categories.tolist()
                median_ordering = df.groupby(by=x)[y].median().sort_values().index.tolist()
                df[x] = df[x].cat.reorder_categories(median_ordering)
                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue)

                # Revert to categorical ordering
                df[x] = df[x].cat.reorder_categories(cat_ordering)
        elif (x in self.cat_features):
            if (y in self.num_features):
                # Sort according to median of numeric
                cat_ordering = df[x].cat.categories.tolist()
                median_ordering = df.groupby(by=x)[y].median().sort_values().index.tolist()
                df[x] = df[x].cat.reorder_categories(median_ordering)

                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue, agg_fcn=agg_fcn)

                # Revert to categorical ordering
                df[x] = df[x].cat.reorder_categories(cat_ordering)
            elif (y in self.pnum_features):
                boxen_and_pie_plot(data=df, x=x, y=y, hue=hue, agg_fcn=agg_fcn)

            elif (y in self.cat_features):
                x_noise, y_noise = np.random.random(len(df))/2, np.random.random(len(df))/2
                self.reset_figure(ncols=2)
                sns.histplot(data=df, x=x, hue=y, stat='count', multiple='stack', ax=self.ax[0])
                self.ax[0].tick_params(axis='x', rotation=90)
                df[x].value_counts(dropna=False).sort_index().plot.pie(autopct='%.0f%%', ax=self.ax[1])

        self.canvas.draw()
        return

    def plot_correlations(self):
        '''Plot largest correlations against target feature.'''
        valid_features = self.get_valid_features()
        prefix, target = self.target_combobox.currentText().split(' ')
        if target == 'none':
            qtw.QMessageBox.warning(
                self, 'No Target Specified', 'Correlations must be computed against some target feature.'
            )
            return

        # Reset canvas, create figure, draw to canvas
        self.reset_figure(ncols=1)     
        (
            self.df[valid_features].corrwith(self.df[target], numeric_only=True)
            .abs().sort_values(ascending=False)[1:self.corr_spinbox.value()+1] #skip correlation with self
            .plot.bar(xlabel='Feature', ylabel='Correlation', ax=self.ax)
        )
        self.canvas.draw()

    def plot_nullpct(self):
        '''Plot features with with largest null percentage. Categorical features have a 'Null' category.'''
        valid_features = self.get_valid_features()
        prefix, target = self.target_combobox.currentText().split(' ')

        # Reset canvas, create figure, draw to canvas
        self.reset_figure(ncols=1)
        (  
            (self.df[valid_features].isna().sum() + self.df.apply(lambda x: sum(x=='Null')))
            .sort_values(ascending=False).divide(len(self.df))[:self.nullpct_spinbox.value()]
            .plot.bar(xlabel='Feature', ylabel='Null Percentage', ax=self.ax)
        )
        self.canvas.draw()

    def encode_transform_feature(self):
        '''Encode a categorical feature from a specified mapping.'''
        # Inputs
        df = self.df
        map_str = self.encode_lineedit.text()
        feature = self.encode_combobox.currentText()
        if feature == 'none':
            return

        try:
            # Create encoding dictionary from input string
            encode_dict = {}
            for entry in map_str.split(','):
                key, value = entry.split('=')
                if key.strip() not in df[feature].values:
                    raise ValueError(f'Key "{key}" is invalid for given feature')
                encode_dict[key.strip()] = float(value.strip())
            
            # Apply encoding
            df['enc_'+feature] = df[feature].map(encode_dict).astype('float')
            self.write_to_feature_engineering( 
                f'df["enc_{feature}"] = df["{feature}"].map({encode_dict}).astype("float")'
            )

            # Add to pnum features
            self.pnum_features.loc['enc_'+feature] = 'enc_'+feature

        # Return exception in messagebox
        except Exception as e:
            qtw.QMessageBox.critical(self, 'Code Error', f'Error message returned: {e}')
            return    

        # Updates
        self.save()
        self.update_menus()
        self.encode_lineedit.setText('')
        self.encode_combobox.setCurrentText('none')
        return

    def log_transform_feature(self):
        '''Apply np.log1p transform to feature.'''
        feature = self.log_combobox.currentText()
        if feature == 'none':
            return

        # Apply log transform
        self.df['log'+feature] = np.log1p(self.df[feature])
        self.write_to_feature_engineering( 
            f'df["log{feature}"]=np.log1p(df["{feature}"])'
        )

        # Add to numerical features
        self.num_features.loc['log'+feature] = 'log'+feature

        # Update
        self.save()
        self.update_menus()
        self.log_combobox.setCurrentText('none')

    def bin_transform_feature(self):
        '''Bin feature into 10 quartiles.'''
        feature = self.bin_combobox.currentText()
        if feature == 'none':
            return

        # Apply binning
        self.df['bin'+feature] = pd.qcut(self.df[feature], 10)
        self.write_to_feature_engineering( 
            f'df["bin{feature}"]=pd.qcut(df["{feature}"], 10)'
        )

        # Add to categorical features
        self.cat_features.loc['bin'+feature] = 'bin'+feature

        # Update
        self.save()
        self.update_menus()
        self.bin_combobox.setCurrentText('none')

    def dummies_transform_feature(self):
        '''Creates dummy encodings for a categorical feature'''
        feature = self.dummies_combobox.currentText()
        if feature == 'none':
            return

        # Apply encoding
        self.df = pd.concat([self.df, pd.get_dummies(self.df[feature], prefix=feature)], axis=1)
        self.write_to_feature_engineering( 
            f'df=pd.concat([df, pd.get_dummies(df["{feature}"], prefix="{feature}")], axis=1)' 
        )

        # Each category encoding to pnum. Drop 'Null' category
        for cat in self.df[feature].unique():
            if cat == 'Null':
                self.df = self.df.drop(columns=[feature+'_Null'])
                self.write_to_feature_engineering( 
                    f'df=df.drop(columns=["{feature}_Null"])'
                )
            else:
                self.pnum_features.loc[feature + '_' + cat] = feature + '_' + cat
        
        # Update
        self.save()
        self.update_menus()
        self.dummies_combobox.setCurrentText('none')

    def drop_feature(self):
        '''Drop feature from consideration. It remains in the dataframe but is hidden from view.'''
        prefix, feature = self.drop_combobox.currentText().split(' ')
        if feature=='none':
            return

        # Add to dropped features
        self.dropped_features.extend([feature])

        # Remove from num, pnum, or cat features
        if prefix=='num':
            self.num_features = self.num_features.drop(feature)
        elif prefix=='pnum':
            self.pnum_features = self.pnum_features.drop(feature)
        elif prefix=='cat':
            self.cat_features = self.cat_features.drop(feature)

        # Updates
        self.save()
        self.update_menus()
        return 

    def undrop_feature(self):
        '''Undrop feature, making it visible in menus again.'''
        feature = self.undrop_combobox.currentText()

        # Remove from dropped features
        self.dropped_features.remove(feature)

        # Add to num, pnum, or cat features
        if is_numeric_dtype(self.df[feature]):
            if self.df[feature].nunique()<11:
                self.pnum_features.loc[feature] = feature
            else:
                self.num_features.loc[feature] = feature
        else:
            self.cat_features.loc[feature] = feature

        # Updates
        self.save()
        self.update_menus()
        return 

    def create_feature_from_code(self):
        '''Run specified code to create new feature(s).'''
        # Inputes
        df = self.df
        old_features = df.columns
        python_code = self.create_lineedit.text()
        

        try:
            exec(python_code)
            self.write_to_feature_engineering(python_code)
        
        # Return exception in messagebox
        except Exception as e:
            qtw.QMessageBox.critical(self, 'Code Error', f'Error message returned: {e}')
            return            
        
        # For each new feature, add to pnum, num, or cat
        for feature in (set(df.columns)-set(old_features)):
            if df[feature].dtype=='object':
                self.cat_features.loc[feature] = feature
                df[feature] = self.df[feature].fillna('Null') #Make null own category
                df[feature] = pd.Categorical(self.df[feature], np.sort(df[feature].unique()))
            else:
                if df[feature].nunique()<11:
                    self.pnum_features.loc[feature] = feature
                else:
                    self.num_features.loc[feature] = feature

        # Updates 
        self.create_lineedit.setText('')
        self.save()
        self.update_menus()

    def write_to_feature_engineering(self, code):
        '''Append code for feature engineering to text file to be implemented later in model_analysis_gui pipeline'''
        filename = self.save_lineedit.text().removesuffix('.csv') + '_feature_engineering.txt'
        if not os.path.exists(filename):
            open(filename, 'w').close()

        with open(filename, 'a') as f:
            f.write(code)
            f.write('\n')
        return

    def load_pptx(self):
        '''Open and return powerpoint slides for notes'''
        filename = self.save_lineedit.text().removesuffix('.csv') + '_slides.pptx'
        ppt = Presentation(filename) if os.path.exists(filename) else Presentation()
        return ppt, filename

    def save_notes_to_pptx(self):
        '''Save comments and figure to slides'''
        # Read inputs
        png_file = 'temp.png'
        title = self.title_lineedit.text()
        text = self.text_textedit.toPlainText()
        self.canvas.print_figure(png_file, bbox_inches='tight')

        # Add slide with text and picture
        ppt, fullpath = self.load_pptx()
        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
        slide.shapes[0].text_frame.paragraphs[0].text = ' ' if title=='' else title
        pic = slide.placeholders[1].insert_picture(png_file)
        pic.crop_top, pic.crop_left, pic.crop_bottom, pic.crop_right = 0,0,0,0
        slide.shapes[2].text_frame.paragraphs[0].text = text
        ppt.save(fullpath)
        os.remove(png_file)

        # Clear text inputs
        self.title_lineedit.setText('')
        self.text_textedit.setText('')

    def get_valid_features(self, prefix=False):
        '''Get valid (undropped) features with optional descriptive prefix.'''
        if prefix:
            valid_features = pd.concat([ 
                'num '+self.num_features,
                'pnum '+self.pnum_features,
                'cat '+self.cat_features
            ]).sort_values().to_list()
        else:
            valid_features = pd.concat([ 
                self.num_features,
                self.pnum_features,
                self.cat_features
            ]).sort_values().to_list()
        return valid_features

    def update_menus(self):
        # Ensure features are sorted. 
        self.dropped_features.sort()
        self.num_features.sort_values()

        num_features_plus_none = self.num_features.to_list()
        num_features_plus_none.insert(0, 'none')

        cat_features_plus_none = self.cat_features.to_list()
        cat_features_plus_none.insert(0, 'none')

        # Combine all valid features with identifying prefix
        valid_features = self.get_valid_features(prefix=True)
        valid_features.insert(0, ' none')


        # Specify options for each combobox
        options_dict = { 
            'x':valid_features,
            'y':valid_features,
            'hue':valid_features,
            'size':valid_features,
            'style':valid_features,
            'target':valid_features,
            'log':num_features_plus_none,
            'bin':num_features_plus_none,
            'dummies':cat_features_plus_none,
            'encode':cat_features_plus_none,
            'drop':valid_features,
            'undrop':self.dropped_features
        }

        # Update combobox items
        for menu, options in options_dict.items():
            combobox = getattr(self, menu+'_combobox')
            current_text = combobox.currentText()
            combobox.clear()
            combobox.addItems(options)
            combobox.setCurrentText(current_text)

        return





if __name__ =='__main__':
    app = qtw.QApplication(sys.argv)
    mw = Main_Window()
    app.exec()


# %%
